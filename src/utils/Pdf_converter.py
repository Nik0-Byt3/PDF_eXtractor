import os
import glob
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import logging
import re

logger = logging.getLogger(__name__)

class BookMarkdownToPPTX:
    def __init__(self, base_path, output_dir):
        self.base_path = base_path
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)
        self.pptx_path = os.path.join(self.output_dir, "book_summary.pptx")
        self.max_chars_per_slide = 600

    def convert_all_chapters(self):
        try:
            logger.info(f"Searching for chapters in {self.base_path}")
            presentation = Presentation()

            chapter_paths = glob.glob(os.path.join(self.base_path, "*_*"))
            if not chapter_paths:
                logger.warning("No chapter directories found.")
                return

            def extract_index(path):
                folder_name = os.path.basename(path)
                match = re.match(r"(\d+)__", folder_name)
                return int(match.group(1)) if match else 9999

            chapter_paths = sorted(chapter_paths, key=extract_index)

            for chapter_path in chapter_paths:
                chapter_name = os.path.basename(chapter_path)
                chapter_title = chapter_name.replace("_", " ")

                summary_path = os.path.join(chapter_path, "Unknown", "unified_summary.md")

                with open(summary_path, "r", encoding="utf-8") as f:
                    content = f.read().strip()

                if not content:
                    logger.warning(f"Empty summary in {chapter_name}")
                    continue

                self._add_title_slide(presentation, chapter_title)

                content_blocks = self._split_content(content, self.max_chars_per_slide)
                for block in content_blocks:
                    self._add_content_slide(presentation, block)

                image_paths = self._find_images(chapter_path)
                for img_path in image_paths:
                    self._add_image_slide(presentation, img_path)

            presentation.save(self.pptx_path)
            logger.info(f"Presentation saved to: {self.pptx_path}")

        except Exception as e:
            logger.error(f"Conversion failed: {str(e)}")

    def _add_title_slide(self, presentation, chapter_title):
        slide_layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title

        title_shape.text = chapter_title
        p = title_shape.text_frame.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.alignment = PP_ALIGN.CENTER

        text_frame = title_shape.text_frame
        slide_height = presentation.slide_height

        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.height = slide_height

    def _add_content_slide(self, presentation, content):
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)

        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title_placeholder.text = ""
        title_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        content_clean = re.sub(r"^#+\s*", "", content, flags=re.MULTILINE).strip()

        text_frame = content_placeholder.text_frame
        text_frame.clear()

        paragraphs = re.split(r'\n{2,}', content_clean)
        for i, para in enumerate(paragraphs):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = para.strip()
            p.font.size = Pt(24)
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(10)
        text_frame.word_wrap = True

        content_placeholder.left = Inches(0.5)
        content_placeholder.top = Inches(1.5)
        content_placeholder.width = Inches(9)
        content_placeholder.height = Inches(5)

    def _add_image_slide(self, presentation, image_path):
        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        slide.shapes.add_picture(image_path, left, top, width=width)

    def _split_content(self, content, max_chars):
        content = content.strip()
        if len(content) <= max_chars:
            return [content]

        blocks = []
        paragraphs = re.split(r'\n{2,}', content)
        current_block = ""

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            if len(current_block) + len(para) + 2 > max_chars:
                if current_block:
                    blocks.append(current_block.strip())
                    current_block = para
                else:
                    for i in range(0, len(para), max_chars):
                        blocks.append(para[i:i+max_chars])
                    current_block = ""
            else:
                current_block += "\n\n" + para if current_block else para

        if current_block:
            blocks.append(current_block.strip())

        return blocks

    def _find_images(self, chapter_path):
        exts = ["*.png", "*.jpg", "*.jpeg", "*.bmp", "*.gif"]
        images = []
        for ext in exts:
            images.extend(glob.glob(os.path.join(chapter_path, "**", ext), recursive=True))
        return sorted(images)

